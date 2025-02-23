import os
import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from email.mime.base import MIMEBase
from email import encoders
from selenium import webdriver
from selenium.webdriver.chrome.options import Options
from selenium.webdriver.common.by import By
from selenium.webdriver.common.keys import Keys
from pptx import Presentation
from pptx.util import Inches
import time
import logging

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

def send_email(subject, body, to_email, attachment_path=None):
    from_email = ""
    from_password = ""
    
    message = MIMEMultipart()
    message["From"] = from_email
    message["To"] = to_email
    message["Subject"] = subject
    message.attach(MIMEText(body, "plain"))
    
    if attachment_path:
        with open(attachment_path, 'rb') as attachment:
            part = MIMEBase('application', 'octet-stream')
            part.set_payload(attachment.read())
            encoders.encode_base64(part)
            part.add_header(
                'Content-Disposition',
                f'attachment; filename={os.path.basename(attachment_path)}'
            )
            message.attach(part)
    
    try:
        server = smtplib.SMTP("smtp.office365.com", 587)
        server.starttls()
        server.login(from_email, from_password)
        server.sendmail(from_email, to_email, message.as_string())
        server.close()
        logger.info("Email sent successfully!")
    except Exception as e:
        logger.error(f"Error sending email: {e}")

def create_powerpoint(screenshot_paths, template_path, output_path):
    try:
        # Load the template
        prs = Presentation(template_path)
        
        # Skip first slide (intro) and add screenshots to slides 2 and 3
        for i, screenshot_path in enumerate(screenshot_paths):
            if i < 2:  # Only process first two screenshots
                slide = prs.slides[i + 1]  # +1 to skip intro slide
                
                # Add screenshot to the empty slide
                # Center the screenshot on the slide
                left = Inches(2)  # 1 inch margin from left
                top = Inches(2)   # 1 inch margin from top
                width = Inches(7)  # 8 inches wide
                height = Inches(4)  # 5.5 inches high
                
                slide.shapes.add_picture(
                    screenshot_path,
                    left,
                    top,
                    width,
                    height
                )
                logger.info(f"Added screenshot to slide {i + 2}")
        
        # Save the presentation
        prs.save(output_path)
        logger.info(f"PowerPoint saved: {output_path}")
        return True
    except Exception as e:
        logger.error(f"PowerPoint creation failed: {e}")
        return False

def main():
    chrome_options = Options()
    chrome_options.add_argument("--headless")
    chrome_options.add_argument("--disable-gpu")
    chrome_options.add_argument("--no-sandbox")
    chrome_options.add_argument("--disable-dev-shm-usage")
    
    driver = webdriver.Chrome(options=chrome_options)
    screenshot_paths = []
    
    logger.info(f"Current working directory: {os.getcwd()}")
    
    try:
        login_url = "http://localhost:5601"
        driver.get(login_url)
        time.sleep(5)
        
        logger.info('Started logging in...')
        username = driver.find_element(By.NAME, "username")
        password = driver.find_element(By.NAME, "password")
        username.send_keys("elastic")
        password.send_keys("kibana123")
        password.send_keys(Keys.RETURN)
        logger.info('Logged in....')
        time.sleep(10)
        
        dashboard_urls = [
            "http://localhost:5601/app/dashboards#/view/722b74f0-b882-11e8-a6d9-e546fe2bba5f?_g=(filters:!())",
            "http://localhost:5601/app/dashboards#/view/7adfa750-4c81-11e8-b3d7-01146121b73d?_g=(filters:!())"
        ]
        
        for i, url in enumerate(dashboard_urls, start=1):
            logger.info('Taking screenshots...')
            driver.get(url)
            time.sleep(15)
            
            total_width = driver.execute_script("return document.body.scrollWidth")
            total_height = driver.execute_script("return document.body.scrollHeight")
            driver.set_window_size(total_width, total_height)
            logger.info(f'H: {total_height} and W: {total_width}')
            
            # screenshot_directory = "/app/screenshots"
            screenshot_directory = "/Users/app/poc/logstashexporter/screenshots"

            if not os.path.exists(screenshot_directory):
                os.makedirs(screenshot_directory)
            
            screenshot_path = os.path.join(screenshot_directory, f"kibana_dashboard_{i}.png")
            driver.get_screenshot_as_file(screenshot_path)
            screenshot_paths.append(screenshot_path)
            logger.info(f"Screenshot saved: {screenshot_path}")
        
        # Create PowerPoint presentation
        # template_path = '/app/screenshots/icrms.pptx'
        template_path = "/Users/app/poc/logstashexporter/icrms.pptx"
        # output_path = "/app/screenshots/template.pptx"
        output_path = "/Users/app/poc/logstashexporter/template.pptx"
        if create_powerpoint(screenshot_paths, template_path, output_path):
            # Send email with PowerPoint attachment
            subject = "Kibana Dashboard Report"
            body = "Please find attached the Kibana dashboard report with screenshots."
            send_email(subject, body, "aflavian@srxconsultant.com", output_path)
            logger.info('email sent..')
        
    finally:
        driver.quit()
        # Clean up screenshots after sending
        for path in screenshot_paths:
            if os.path.exists(path):
                os.remove(path)
        logger.info('Cleaned up the screenshots')


if __name__ == "__main__":
    main()